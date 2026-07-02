from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

template_path = os.path.join(os.path.dirname(__file__), "template.pptx")
output_path = os.path.join(os.path.dirname(__file__), "Decision_Intelligence_Platform.pptx")

# Open template to inherit its theme
prs = Presentation(template_path)

# Remove all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Get the first slide layout from the template
slide_layouts = prs.slide_layouts
title_layout = None
blank_layout = None
for layout in slide_layouts:
    if layout.name == "Title Slide" or "title" in layout.name.lower():
        title_layout = layout
    if layout.name == "Blank" or "blank" in layout.name.lower():
        blank_layout = layout

if not title_layout:
    title_layout = slide_layouts[0]
if not blank_layout:
    blank_layout = slide_layouts[-1]

def add_slide(layout=None):
    if layout is None:
        layout = blank_layout
    slide = prs.slides.add_slide(layout)
    # Remove all placeholders
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return slide

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=RGBColor(0xFF, 0xFF, 0xFF), bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=RGBColor(0xCC, 0xCC, 0xCC), spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
    return txBox

# Colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2 = RGBColor(0xFB, 0x92, 0x3C)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)

W = Inches(13.33)
H = Inches(7.5)

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = add_slide(title_layout)
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
             "Decision Intelligence Platform", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "AI-Powered Insights for Smarter Community Decisions", 24, GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.8),
             "Google Solution Challenge 2026", 20, ACCENT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "THE PROBLEM", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
             "Why communities struggle with decision-making", 16, GRAY)

problems = [
    "Communities generate massive amounts of data from transport, healthcare,\n    utilities, citizen feedback — but it's siloed across departments",
    "Transforming raw data into actionable insights is slow and complex,\n    leading to reactive rather than predictive decisions",
    "No unified platform exists that connects AI, analytics, and\n    decision support for community stakeholders",
    "Citizen engagement and community well-being suffer without\n    data-driven governance and transparent decision-making",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), problems, 14, GRAY, Pt(10))

# Right box
add_text_box(slide, Inches(7.5), Inches(0.5), Inches(5), Inches(0.7), "THE IMPACT", 28, ACCENT2, True)
box = add_shape(slide, Inches(7.5), Inches(1.4), Inches(5), Inches(5.4), fill_color=CARD_BG)
metrics = [
    "70% of municipal data is unstructured and underutilized",
    "Decision-making cycles average 3-6 months in most cities",
    "Citizen satisfaction drops 15% with slow service responses",
    "Predictive analytics capabilities absent in 80% of communities",
    "AI adoption in public sector lags 5+ years behind private sector",
]
add_bullet_slide(slide, Inches(7.8), Inches(1.6), Inches(4.5), Inches(5), metrics, 14, ACCENT2, Pt(8))

# ============================================================
# SLIDE 3: Our Solution
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "OUR SOLUTION", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Decision Intelligence Platform — AI-powered, data-driven, action-oriented", 14, GRAY)

features = [
    ("Conversational AI Assistant", "Natural language interface powered by Google Gemini LLM for querying any domain"),
    ("Multi-Domain Analytics", "8 integrated domains with dashboards, trends, and performance metrics"),
    ("Predictive Forecasting", "ML-powered predictions for traffic, energy, health, and environmental outcomes"),
    ("Decision Intelligence Engine", "Multi-criteria analysis that evaluates options and recommends optimal actions"),
    ("Intelligent Data Ingestion", "Connect CSV, JSON, PDF, APIs, and webhooks — auto-categorized by domain"),
    ("Responsible & Explainable AI", "Every insight includes confidence scores, source citations, and reasoning"),
]
y = 1.6
for title, desc in features:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.35), f"▸ {title}", 18, ACCENT, True)
    add_text_box(slide, Inches(1.2), Inches(y + 0.35), Inches(10.5), Inches(0.3), desc, 13, GRAY)
    y += 0.85

# ============================================================
# SLIDE 4: Architecture
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "SYSTEM ARCHITECTURE", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Built on Google Cloud ecosystem with modern full-stack design", 14, GRAY)

layers = [
    ("FRONTEND", "Next.js 15 · TypeScript · Tailwind CSS · Recharts",
     "Dashboard · AI Chat · Analytics · Decision Tools · Data Management"),
    ("API GATEWAY", "FastAPI · RESTful · WebSockets · Async · Auto-docs",
     "Unified API layer for all platform services with OpenAPI documentation"),
    ("AI SERVICES", "Google Gemini LLM · Vertex AI · RAG Engine",
     "Natural language understanding · Context-aware responses · Domain expertise"),
    ("DATA LAYER", "Sample Data Generators · ChromaDB · BigQuery (optional)",
     "8 domain datasets with 82000+ records · Vector search · Real-time analytics"),
    ("INFRASTRUCTURE", "Docker · Google Cloud Run · Cloud Build CI/CD",
     "Containerized microservices · Auto-scaling · Production-ready deployment"),
]
y = 1.5
for title, tech, desc in layers:
    box = add_shape(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.95), fill_color=CARD_BG)
    add_text_box(slide, Inches(1), Inches(y + 0.05), Inches(3), Inches(0.35), title, 14, ACCENT, True)
    add_text_box(slide, Inches(4.2), Inches(y + 0.05), Inches(8), Inches(0.35), tech, 12, GRAY)
    add_text_box(slide, Inches(1), Inches(y + 0.4), Inches(11), Inches(0.4), desc, 11, WHITE)
    y += 1.1

# ============================================================
# SLIDE 5: 8 Domains
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "8 SOLUTION DOMAINS", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Comprehensive coverage of community decision-making needs", 14, GRAY)

domains = [
    ("Urban Mobility", "Traffic optimization, transit planning,\ncongestion prediction"),
    ("Public Safety", "Emergency response, crime analytics,\ndisaster preparedness"),
    ("Healthcare", "Community wellness, resource allocation,\naccess optimization"),
    ("Environment", "Air/water quality, waste management,\nsustainability tracking"),
    ("Energy", "Smart grid, consumption analysis,\nefficiency recommendations"),
    ("Education", "Student performance, resource planning,\nprogram effectiveness"),
    ("Citizen Engagement", "Feedback analysis, service optimization,\nparticipation tracking"),
    ("Tourism & Economy", "Visitor patterns, local spending,\neconomic development"),
]
x_start, y_start = Inches(0.8), Inches(1.6)
for i, (title, desc) in enumerate(domains):
    col = i % 4
    row = i // 4
    cx = x_start + col * Inches(3.0)
    cy = y_start + row * Inches(2.6)
    box = add_shape(slide, cx, cy, Inches(2.7), Inches(2.2), fill_color=CARD_BG)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.15),
                 Inches(2.4), Inches(0.4), title, 16, ACCENT, True)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.6),
                 Inches(2.4), Inches(1.4), desc, 12, GRAY)

# ============================================================
# SLIDE 6: Tech Stack
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "TECHNOLOGY STACK", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
             "Google Cloud ecosystem + modern open-source technologies", 14, GRAY)

techs = [
    ("Google Gemini", "Large Language Model for NLU and generation"),
    ("Vertex AI", "ML model training & deployment pipeline"),
    ("BigQuery", "Serverless data warehouse for analytics"),
    ("Cloud Run", "Managed container platform, auto-scaling"),
    ("FastAPI", "High-performance async Python API framework"),
    ("Next.js 15", "React framework with SSR & app router"),
    ("ChromaDB", "Vector database for RAG & semantic search"),
    ("Docker/Cloud Build", "Containerization & CI/CD pipeline"),
]
x_start, y_start = Inches(0.8), Inches(1.6)
for i, (title, desc) in enumerate(techs):
    col = i % 4
    row = i // 4
    cx = x_start + col * Inches(3.0)
    cy = y_start + row * Inches(2.6)
    box = add_shape(slide, cx, cy, Inches(2.7), Inches(2.1), fill_color=CARD_BG)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.15),
                 Inches(2.4), Inches(0.4), title, 16, ACCENT, True)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.6),
                 Inches(2.4), Inches(1.3), desc, 12, GRAY)

# ============================================================
# SLIDE 7: Demo Flow
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "DEMONSTRATION WALKTHROUGH", 36, WHITE, True)

steps = [
    "1. Dashboard Overview — Cross-domain performance metrics, activity feed, AI insights summary",
    "2. AI Chat — Ask questions in natural language; get data-driven answers with sources & suggestions",
    "3. Domain Analytics — Select any domain; view trends, forecasts, and actionable insights",
    "4. Decision Intelligence — Input problem context & options; AI evaluates & recommends best action",
    "5. Data Source Management — Connect CSV, JSON, PDF, API, webhook; auto-categorized by domain",
    "6. Smart Forecasts — Generate predictive projections across all domains with confidence scores",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5), steps, 18, GRAY, Pt(14))

# ============================================================
# SLIDE 8: Expected Impact
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "EXPECTED IMPACT", 36, WHITE, True)

impacts = [
    ("65%", "faster decision-making\nwith AI-powered insights"),
    ("40%", "improvement in\nresource allocation efficiency"),
    ("30%", "reduction in response\ntimes for critical services"),
    ("25%", "increase in citizen\nsatisfaction scores"),
    ("50%", "of data processing\nfully automated"),
    ("8", "integrated community\ndomains covered"),
]
x_start, y_start = Inches(0.8), Inches(1.8)
for i, (num, label) in enumerate(impacts):
    col = i % 3
    row = i // 3
    cx = x_start + col * Inches(4.0)
    cy = y_start + row * Inches(2.6)
    box = add_shape(slide, cx, cy, Inches(3.5), Inches(2.2), fill_color=CARD_BG)
    add_text_box(slide, cx + Inches(0.3), cy + Inches(0.2),
                 Inches(3), Inches(0.7), num, 44, ACCENT, True, PP_ALIGN.CENTER)
    add_text_box(slide, cx + Inches(0.3), cy + Inches(1.0),
                 Inches(3), Inches(1.0), label, 14, GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: Future Roadmap
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), "FUTURE ROADMAP", 36, WHITE, True)

roadmap = [
    "Phase 1: Real-time data pipeline with Google Pub/Sub and Dataflow streaming",
    "Phase 2: Multi-modal AI — image/video understanding for traffic, safety, environment",
    "Phase 3: Community mobile app with push notifications and citizen reporting",
    "Phase 4: Custom ML models via Vertex AI AutoML trained on community-specific data",
    "Phase 5: Cross-community benchmarking and collaborative decision-making dashboards",
    "Phase 6: Integration with Google Maps API & BigQuery GIS for spatial analytics",
    "Phase 7: IoT sensor network integration for real-time environmental monitoring",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5), roadmap, 16, GRAY, Pt(10))

# ============================================================
# SLIDE 10: Thank You
# ============================================================
slide = add_slide(title_layout)
add_shape(slide, 0, 0, W, H, fill_color=DARK_BG)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
             "Thank You", 56, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
             "Decision Intelligence Platform", 26, ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.1), Inches(11), Inches(0.5),
             "AI-Powered Insights for Smarter Community Decisions", 18, GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "GitHub: github.com/kushaljagad73/decision-intelligence-platform", 14, GRAY, alignment=PP_ALIGN.CENTER)

prs.save(output_path)
print(f"PPT saved to: {output_path}")
