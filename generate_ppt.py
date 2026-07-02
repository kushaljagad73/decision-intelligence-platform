from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1E, 0x40, 0xAF)
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
    return txBox


# SLIDE 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
             "Decision Intelligence Platform", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "AI-Powered Insights for Smarter Community Decisions", 24, GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.8),
             "Google Solution Challenge 2026", 20, ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
             "Leveraging AI, Data Analytics & Intelligent Automation", 16, GRAY, alignment=PP_ALIGN.CENTER)

# SLIDE 2: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "THE PROBLEM", 32, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
             "Why communities struggle with decision-making", 16, GRAY)

problems = [
    "Communities generate massive data from transport, healthcare, utilities, citizen feedback",
    "Information is siloed across departments - no unified view",
    "Transforming raw data into actionable insights is complex and slow",
    "Decisions are often reactive rather than predictive",
    "Citizen engagement and community well-being suffer without data-driven governance",
    "No single platform connects AI, analytics, and decision support",
]
add_bullet_slide(slide, Inches(0.8), Inches(2), Inches(5.5), Inches(5), problems, 15, GRAY)

# Right side metrics box
add_text_box(slide, Inches(7.5), Inches(0.5), Inches(5), Inches(0.7), "IMPACT", 28, ACCENT, True)

metrics_box = slide.shapes.add_shape(
    1, Inches(7.5), Inches(1.4), Inches(5), Inches(5.4))
metrics_box.fill.solid()
metrics_box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
metrics_box.line.fill.background()

metrics = [
    "70% of municipal data is unstructured and underutilized",
    "Decision-making cycles average 3-6 months",
    "Citizen satisfaction drops 15% with slow responses",
    "Predictive capabilities are absent in 80% of communities",
    "AI adoption in public sector lags by 5+ years",
]
add_bullet_slide(slide, Inches(7.8), Inches(1.6), Inches(4.5), Inches(5), metrics, 15, RGBColor(0xFB, 0x92, 0x3C))

# SLIDE 3: Our Solution
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "OUR SOLUTION", 32, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "Decision Intelligence Platform - AI-powered, data-driven, action-oriented", 16, GRAY)

# Key features
features = [
    ("Conversational AI Assistant", "Natural language interface powered by Google Gemini LLM for querying any domain data"),
    ("Multi-Domain Analytics", "8 integrated domains with real-time dashboards, trends, and performance metrics"),
    ("Predictive Forecasting", "ML-powered predictions for traffic, energy, health, and environmental outcomes"),
    ("Decision Intelligence Engine", "Multi-criteria analysis that evaluates options and recommends optimal actions"),
    ("Intelligent Data Ingestion", "Connect CSV, JSON, PDF, APIs, and webhooks - auto-categorized by domain"),
    ("Responsible & Explainable AI", "Every insight includes confidence scores, source citations, and reasoning"),
]
y = 2.0
for title, desc in features:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.4), title, 18, ACCENT, True)
    add_text_box(slide, Inches(0.8), Inches(y + 0.4), Inches(11), Inches(0.3), desc, 14, GRAY)
    y += 0.85

# SLIDE 4: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "SYSTEM ARCHITECTURE", 32, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "Built on Google Cloud ecosystem with modern full-stack design", 16, GRAY)

# Architecture layers
layers = [
    ("FRONTEND", "Next.js 15 · TypeScript · Tailwind CSS · Recharts",
     "Dashboard · Chat Interface · Analytics · Decision Tools · Data Management"),
    ("API GATEWAY", "FastAPI · RESTful · WebSockets · CORS · Async",
     "Unified API for all platform services with auto-generated OpenAPI docs"),
    ("AI SERVICES", "Google Gemini LLM · Vertex AI · RAG Engine",
     "Natural language understanding · Context-aware responses · Domain expertise"),
    ("DATA LAYER", "Sample Data Generators · ChromaDB · BigQuery (opt)",
     "8 domain datasets with 6000+ records · Vector search · Real-time analytics"),
    ("INFRASTRUCTURE", "Docker · Google Cloud Run · Cloud Build",
     "Containerized microservices · Auto-scaling · CI/CD pipeline"),
]
y = 1.9
for title, tech, desc in layers:
    box = slide.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(11.5), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    box.line.fill.background()
    add_text_box(slide, Inches(1), Inches(y + 0.05), Inches(3), Inches(0.35), title, 14, ACCENT, True)
    add_text_box(slide, Inches(4.2), Inches(y + 0.05), Inches(8), Inches(0.35), tech, 13, GRAY)
    add_text_box(slide, Inches(1), Inches(y + 0.4), Inches(11), Inches(0.4), desc, 12, WHITE)
    y += 1.05

# SLIDE 5: 8 Domains
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "8 SOLUTION DOMAINS", 32, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "Comprehensive coverage of community decision-making needs", 16, GRAY)

domains = [
    ("Urban Mobility", "Traffic optimization, transit planning,\ncongestion prediction"),
    ("Public Safety", "Emergency response, crime analytics,\ndisaster preparedness"),
    ("Healthcare", "Community wellness, resource allocation,\naccess optimization"),
    ("Environment", "Air/water quality, waste management,\nsustainability tracking"),
    ("Energy", "Smart grid, consumption analysis,\nefficiency recommendations"),
    ("Education", "Student performance, resource planning,\nprogram effectiveness"),
    ("Citizen Engagement", "Feedback analysis, service optimization,\nparticipation tracking"),
    ("Tourism & Economy", "Visitor patterns, local spending,\neconomic development"),
]
x, y = 0.8, 2.0
for i, (title, desc) in enumerate(domains):
    col = i % 4
    row = i // 4
    cx = Inches(x + col * 3.0)
    cy = Inches(y + row * 2.4)

    box = slide.shapes.add_shape(1, cx, cy, Inches(2.7), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    box.line.fill.background()

    add_text_box(slide, Inches(cx.inches + 0.15), Inches(cy.inches + 0.15),
                 Inches(2.4), Inches(0.4), title, 16, ACCENT, True)
    add_text_box(slide, Inches(cx.inches + 0.15), Inches(cy.inches + 0.55),
                 Inches(2.4), Inches(1.3), desc, 12, GRAY)

# SLIDE 6: Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "TECHNOLOGY STACK", 32, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "Google Cloud ecosystem + modern open-source technologies", 16, GRAY)

techs = [
    ("Google Gemini", "Large Language Model for natural language understanding and generation"),
    ("Google Vertex AI", "ML model training, deployment, and MLOps pipeline"),
    ("Google BigQuery", "Serverless data warehouse for large-scale analytics"),
    ("Google Cloud Run", "Fully managed container platform with auto-scaling"),
    ("FastAPI (Python)", "High-performance async API framework with auto-docs"),
    ("Next.js 15", "React framework with SSR, app router, and optimal performance"),
    ("ChromaDB", "Open-source vector database for RAG and semantic search"),
    ("Docker + Cloud Build", "Containerization and CI/CD for reliable deployments"),
]
x, y = 0.8, 2.0
for i, (title, desc) in enumerate(techs):
    col = i % 4
    row = i // 4
    cx = Inches(x + col * 3.0)
    cy = Inches(y + row * 2.4)

    box = slide.shapes.add_shape(1, cx, cy, Inches(2.7), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    box.line.fill.background()

    add_text_box(slide, Inches(cx.inches + 0.15), Inches(cy.inches + 0.15),
                 Inches(2.4), Inches(0.4), title, 16, ACCENT, True)
    add_text_box(slide, Inches(cx.inches + 0.15), Inches(cy.inches + 0.55),
                 Inches(2.4), Inches(1.3), desc, 12, GRAY)

# SLIDE 7: Demo Flow
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "DEMONSTRATION", 32, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "End-to-end platform walkthrough", 16, GRAY)

steps = [
    "1. Dashboard Overview - View cross-domain performance metrics and real-time activity feed",
    "2. AI Chat - Ask 'Analyze downtown traffic patterns' - get data-driven insights with sources",
    "3. Analytics - Select domains, view trends, forecasts, and actionable insights",
    "4. Decision Intelligence - Input options, get AI-recommended actions with pros/cons analysis",
    "5. Data Sources - Connect new data sources, manage ingestion across 8 domains",
    "6. Predictive Forecasting - Generate forward-looking predictions with confidence scores",
]
add_bullet_slide(slide, Inches(0.8), Inches(2), Inches(11), Inches(4.5), steps, 18, GRAY, Pt(14))

# SLIDE 8: Impact
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "EXPECTED IMPACT", 32, WHITE, True)

impacts = [
    ("65%", "faster decision-making\nwith AI-powered insights"),
    ("40%", "improvement in\nresource allocation efficiency"),
    ("30%", "reduction in response\ntimes for critical services"),
    ("25%", "increase in citizen\nsatisfaction scores"),
    ("50%", "of data processing\nfully automated"),
    ("8", "integrated community\ndomains covered"),
]
x, y = 0.8, 2.0
for i, (num, label) in enumerate(impacts):
    col = i % 3
    row = i // 3
    cx = Inches(x + col * 4.0)
    cy = Inches(y + row * 2.4)

    box = slide.shapes.add_shape(1, cx, cy, Inches(3.5), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    box.line.fill.background()

    add_text_box(slide, Inches(cx.inches + 0.3), Inches(cy.inches + 0.2),
                 Inches(3), Inches(0.7), num, 40, ACCENT, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(cx.inches + 0.3), Inches(cy.inches + 1.0),
                 Inches(3), Inches(0.8), label, 14, GRAY, alignment=PP_ALIGN.CENTER)

# SLIDE 9: Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "FUTURE ROADMAP", 32, WHITE, True)

roadmap = [
    "Phase 1: Real-time data pipeline with Google Pub/Sub and Dataflow integration",
    "Phase 2: Multi-modal AI - image/video understanding for traffic, safety, environment",
    "Phase 3: Community mobile app with push notifications and citizen reporting",
    "Phase 4: Advanced ML models - custom Vertex AI training on community-specific data",
    "Phase 5: Cross-community benchmarking and collaborative decision-making",
    "Phase 6: Integration with Google Maps API, BigQuery GIS for spatial analytics",
    "Phase 7: IoT sensor network integration for real-time environmental monitoring",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), roadmap, 16, GRAY, Pt(10))

# SLIDE 10: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
             "Thank You", 56, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "Decision Intelligence Platform", 24, ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
             "AI-Powered Insights for Smarter Community Decisions", 18, GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "GitHub: github.com/kushaljagad73/decision-intelligence-platform", 14, GRAY, alignment=PP_ALIGN.CENTER)

output_path = os.path.join(os.path.dirname(__file__), "Decision_Intelligence_Platform.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
