from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

template_path = os.path.join(os.path.dirname(__file__), "template.pptx")
output_path = os.path.join(os.path.dirname(__file__), "DIP_Prototype_Deck.pptx")

prs = Presentation(template_path)

while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

slide_layouts = prs.slide_layouts
blank_layout = slide_layouts[-1]

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return slide

def add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=RGBColor(0xFF, 0xFF, 0xFF), bold=False,
                 alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = alignment
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=14,
                color=RGBColor(0xCC, 0xCC, 0xCC), spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = spacing
    return txBox

def add_multiline_box(slide, left, top, width, height, text, font_size=13,
                      color=RGBColor(0xCC, 0xCC, 0xCC)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(2)
    return txBox

DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2 = RGBColor(0xFB, 0x92, 0x3C)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xEF, 0x44, 0x44)

W, H = Inches(13.33), Inches(7.5)
M = Inches(0.7)

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(2.0), Inches(11), Inches(1.5),
             "Decision Intelligence Platform", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, M, Inches(3.5), Inches(11), Inches(0.8),
             "AI-Powered Insights for Smarter Community Decisions", 22, GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, M, Inches(5.0), Inches(11), Inches(0.6),
             "Google Solution Challenge 2026", 20, ACCENT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: BRIEF ABOUT IDEA
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "BRIEF ABOUT IDEA", 32, WHITE, True)

left_text = (
    "Modern communities generate vast amounts of structured and unstructured data from\n"
    "transportation systems, healthcare services, environmental monitoring, citizen feedback,\n"
    "utility networks, and public services. However, this data remains siloed and underutilized.\n\n"
    "The Decision Intelligence Platform bridges this gap by providing an AI-powered, unified\n"
    "platform that ingests data from multiple sources, understands natural language queries,\n"
    "generates actionable insights, predicts outcomes, and recommends optimal decisions.\n\n"
    "Built on Google Cloud ecosystem with Gemini LLM at its core, the platform covers 8 critical\n"
    "community domains and empowers stakeholders — from city officials to citizens — to make\n"
    "data-driven decisions that improve quality of life and community well-being."
)
add_multiline_box(slide, M, Inches(1.2), Inches(7.5), Inches(5.5), left_text, 14, GRAY)

# Right panel - Key highlights
add_shape(slide, Inches(8.5), Inches(1.2), Inches(4.3), Inches(5.5), CARD_BG)
highlights = [
    "AI-Powered Conversational Analytics",
    "8 Integrated Community Domains",
    "Multi-Criteria Decision Engine",
    "Predictive Forecasting & Trends",
    "Explainable & Responsible AI",
    "Real-time Data Ingestion",
    "Google Cloud Native",
    "Offline-First Architecture",
]
add_bullets(slide, Inches(8.7), Inches(1.4), Inches(4), Inches(5), highlights, 14, ACCENT, Pt(8))

# ============================================================
# SLIDE 3: SOLUTION CAPABILITIES
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "SOLUTION CAPABILITIES", 32, WHITE, True)
add_text_box(slide, M, Inches(0.9), Inches(10), Inches(0.4),
             "What the Decision Intelligence Platform can do", 14, GRAY)

capabilities = [
    ("Understand & Analyze Data", "Ingests structured/unstructured data from 8 domains, uses Gemini LLM for natural language understanding and contextual analysis"),
    ("Answer Questions in Natural Language", "Conversational AI interface allows users to ask questions like 'Show me traffic trends downtown' and get instant answers"),
    ("Identify Patterns & Anomalies", "Advanced analytics engine detects trends, outliers, and correlations across data sources automatically"),
    ("Generate Recommendations", "Decision intelligence engine uses multi-criteria analysis to evaluate options and recommend optimal actions"),
    ("Automate Workflows", "Intelligent data ingestion pipeline auto-categorizes, processes, and indexes data by domain"),
    ("Support Decision-Making", "Every insight includes confidence scores, source citations, reasoning, and suggested follow-up actions"),
]
y = 1.5
for title, desc in capabilities:
    add_shape(slide, M, Inches(y), Inches(11.5), Inches(0.8), CARD_BG)
    add_text_box(slide, Inches(0.9), Inches(y + 0.05), Inches(4), Inches(0.35), f"▸ {title}", 15, ACCENT, True)
    add_text_box(slide, Inches(0.9), Inches(y + 0.38), Inches(11), Inches(0.4), desc, 12, GRAY)
    y += 0.95

# ============================================================
# SLIDE 4: OPPORTUNITIES
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "OPPORTUNITIES", 32, WHITE, True)

opportunities = [
    ("Urban Mobility & Transportation", "Optimize traffic flow, reduce congestion, improve public transit efficiency through AI-powered route analysis and demand prediction"),
    ("Public Safety & Emergency Response", "Reduce response times with predictive dispatch, identify high-risk areas, and optimize resource allocation for emergencies"),
    ("Healthcare Access & Community Wellness", "Improve healthcare accessibility, predict patient demand, optimize clinic staffing, and identify community health trends"),
    ("Environmental Sustainability & Climate", "Monitor air/water quality in real-time, optimize waste management routes, track carbon emissions, and support green initiatives"),
    ("Energy Efficiency & Smart Utilities", "Reduce energy consumption through smart grid analytics, predict peak demand, and optimize renewable energy integration"),
    ("Citizen Engagement & Public Services", "Analyze feedback sentiment, improve service delivery, increase participation through multi-channel digital engagement platforms"),
]
y = 1.2
for title, desc in opportunities:
    add_shape(slide, M, Inches(y), Inches(11.5), Inches(0.9), CARD_BG)
    add_text_box(slide, Inches(0.9), Inches(y + 0.05), Inches(4.5), Inches(0.35), f"▸ {title}", 15, ACCENT2, True)
    add_text_box(slide, Inches(0.9), Inches(y + 0.4), Inches(11), Inches(0.4), desc, 12, GRAY)
    y += 1.0

# ============================================================
# SLIDE 5: FEATURES
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "FEATURES OFFERED", 32, WHITE, True)

features = [
    ("Conversational AI Assistant", "Natural language interface powered by Gemini LLM for querying any domain with context-aware responses"),
    ("Multi-Domain Analytics Dashboard", "Real-time visualization of KPIs across 8 community domains with drill-down capabilities"),
    ("Predictive Forecasting Engine", "ML-powered projections for traffic, energy demand, health trends, and environmental metrics"),
    ("Decision Intelligence Engine", "Multi-criteria analysis that evaluates options, ranks alternatives, and recommends optimal actions"),
    ("Intelligent Data Ingestion", "Connect CSV, JSON, PDF, APIs, and webhooks with auto-categorization by domain"),
    ("RAG-Powered Knowledge Base", "Retrieval-Augmented Generation using vector search for accurate, source-cited responses"),
    ("Explainable AI Outputs", "Every recommendation includes confidence scores, reasoning, source citations, and trade-off analysis"),
    ("Real-time Activity Monitoring", "Live feed of system events, data ingestion status, insight generation, and decision tracking"),
]
x_start, y_start = M, Inches(1.2)
for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    cx = x_start + col * Inches(6.2)
    cy = y_start + row * Inches(1.45)
    add_shape(slide, cx, cy, Inches(5.8), Inches(1.25), CARD_BG)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.08), Inches(5.5), Inches(0.3), f"✓ {title}", 14, ACCENT, True)
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.4), Inches(5.5), Inches(0.75), desc, 11, GRAY)

# ============================================================
# SLIDE 6: PROCESS FLOW
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "PROCESS FLOW DIAGRAM", 32, WHITE, True)

# Flow: Data Ingestion → Processing → AI Analysis → Insights → Decisions → Action
flow_steps = [
    ("1. DATA INGESTION", "CSV / JSON / PDF\nAPI / Webhook\n8 domain sources", ACCENT),
    ("2. DATA PROCESSING", "Auto-categorization\nVector embedding\nIndexing", GREEN),
    ("3. AI ANALYSIS", "Gemini LLM\nRAG Context\nPattern Detection", ACCENT2),
    ("4. INSIGHT GEN", "Trend Analysis\nPredictive Forecast\nAnomaly Alerts", RED),
    ("5. DECISIONS", "Multi-criteria Eval\nOption Ranking\nPros/Cons", GREEN),
    ("6. ACTION", "Recommendations\nReports\nDashboard Updates", ACCENT),
]

box_w = Inches(1.8)
box_h = Inches(2.2)
arrow_w = Inches(0.4)
start_x = Inches(0.5)
flow_y = Inches(2.0)

for i, (title, desc, color) in enumerate(flow_steps):
    cx = start_x + i * (box_w + Inches(0.35))
    add_shape(slide, cx, flow_y, box_w, box_h, CARD_BG)
    add_text_box(slide, cx + Inches(0.08), flow_y + Inches(0.1), box_w - Inches(0.16), Inches(0.4), title, 10, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, cx + Inches(0.08), flow_y + Inches(0.5), box_w - Inches(0.16), Inches(1.5), desc, 10, GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(flow_steps) - 1:
        arrow_x = cx + box_w + Inches(0.02)
        add_text_box(slide, arrow_x, flow_y + Inches(0.8), arrow_w, Inches(0.4), "→", 24, GRAY, alignment=PP_ALIGN.CENTER)

# Bottom description
add_multiline_box(slide, M, Inches(4.8), Inches(11), Inches(2),
    "The platform follows an end-to-end pipeline: Raw data enters through ingestion connectors → auto-categorized and vectorized → "
    "analyzed by Gemini LLM with RAG context → patterns and forecasts extracted → options evaluated against multiple criteria → "
    "actionable recommendations delivered through dashboard, chat, and reports.\n\n"
    "Feedback loop: User interactions and decisions are logged to improve future recommendations, creating a continuous "
    "learning cycle that gets smarter over time.", 12, GRAY)

# ============================================================
# SLIDE 7: USE-CASE DIAGRAM
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "USE-CASE DIAGRAM", 32, WHITE, True)

# Central actor boxes
actors = [
    ("City Administrators", "Dashboard monitoring\nPolicy decisions\nResource allocation\nPerformance tracking", ACCENT),
    ("Citizens", "Query services\nReport issues\nView analytics\nSubmit feedback", GREEN),
    ("Emergency Services", "Dispatch optimization\nResponse analytics\nRisk assessment\nResource planning", ACCENT2),
    ("Urban Planners", "Traffic analysis\nTransit optimization\nInfrastructure planning\nSustainability goals", RED),
]

for i, (actor, usecases, color) in enumerate(actors):
    cx = M + i * Inches(3.1)
    add_shape(slide, cx, Inches(1.3), Inches(2.8), Inches(1.0), color)
    add_text_box(slide, cx + Inches(0.1), Inches(1.35), Inches(2.6), Inches(0.5), actor, 15, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, cx + Inches(0.1), Inches(1.75), Inches(2.6), Inches(0.4), usecases.split("\n")[0], 10, WHITE, alignment=PP_ALIGN.CENTER)

    add_shape(slide, cx, Inches(2.5), Inches(2.8), Inches(1.6), CARD_BG)
    add_bullets(slide, cx + Inches(0.15), Inches(2.6), Inches(2.5), Inches(1.4),
                usecases.split("\n")[1:], 10, GRAY, Pt(4))

# Central platform
add_shape(slide, Inches(4.5), Inches(4.5), Inches(4.3), Inches(1.5), ACCENT)
add_text_box(slide, Inches(4.5), Inches(4.7), Inches(4.3), Inches(0.5),
             "DECISION INTELLIGENCE PLATFORM", 16, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.6),
             "AI Engine · Analytics · Decisions · Data", 12, WHITE, alignment=PP_ALIGN.CENTER)

# Connecting lines description
add_text_box(slide, Inches(4.5), Inches(6.3), Inches(4.3), Inches(0.4),
             "↑ All actors interact through the unified platform ↑", 11, GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: ARCHITECTURE
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "ARCHITECTURE DIAGRAM", 32, WHITE, True)

layers = [
    ("PRESENTATION LAYER", "Next.js 15 · Tailwind CSS · Recharts",
     "Dashboard | AI Chat | Analytics | Decisions | Data Management | Responsive UI"),
    ("API GATEWAY LAYER", "FastAPI · REST · WebSocket · Async",
     "Authentication | Rate Limiting | Request Validation | API Documentation (OpenAPI)"),
    ("AI SERVICES LAYER", "Google Gemini · RAG Engine · Vector Search",
     "Natural Language Understanding | Context Retrieval | Pattern Detection | Forecasting"),
    ("DATA & STORAGE LAYER", "ChromaDB · BigQuery · File Store",
     "Vector Embeddings | Time-series Data | Sample Generators | 8 Domain Schemas"),
    ("INFRASTRUCTURE LAYER", "Docker · Cloud Run · Cloud Build",
     "Container Orchestration | Auto-scaling | CI/CD Pipeline | Load Balancing"),
]
y = 1.2
for title, tech, desc in layers:
    add_shape(slide, M, Inches(y), Inches(11.5), Inches(1.05), CARD_BG)
    add_text_box(slide, Inches(0.9), Inches(y + 0.05), Inches(3.5), Inches(0.3), title, 14, ACCENT, True)
    add_text_box(slide, Inches(4.5), Inches(y + 0.05), Inches(7.5), Inches(0.3), tech, 12, GRAY)
    add_text_box(slide, Inches(0.9), Inches(y + 0.4), Inches(11), Inches(0.5), desc, 11, WHITE)
    y += 1.15

# System arrows
add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.4), Inches(0.03), ACCENT)
add_shape(slide, Inches(0.8), Inches(6.35), Inches(11.4), Inches(0.03), ACCENT)

# ============================================================
# SLIDE 9: WIREFRAMES
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "WIREFRAMES / MOCK DIAGRAMS", 32, WHITE, True)

screens = [
    ("Dashboard", "Cross-domain KPI cards\nDomain performance grid\nReal-time activity feed\nAI insights summary", ACCENT),
    ("AI Chat Interface", "Conversational UI\nDomain selector sidebar\nMessage history\nSuggestion chips", GREEN),
    ("Analytics View", "Domain filter tabs\nTrend line charts\nKey insights panel\nForecast projections", ACCENT2),
    ("Decision Tool", "Context input area\nOptions list editor\nRecommendation card\nPros/cons comparison", RED),
]

for i, (name, desc, color) in enumerate(screens):
    cx = M + i * Inches(3.1)
    # Screen frame
    add_shape(slide, cx, Inches(1.3), Inches(2.8), Inches(3.5), CARD_BG)
    # Screen header
    add_shape(slide, cx, Inches(1.3), Inches(2.8), Inches(0.5), color)
    add_text_box(slide, cx + Inches(0.1), Inches(1.35), Inches(2.6), Inches(0.4), name, 13, WHITE, True, PP_ALIGN.CENTER)
    # Screen content mockup
    lines = [
        "┌─────────────────┐",
        "│  [Metric Cards]  │",
        "├──────┬──────┬───┤",
        "│ Dom  │ Dom  │Dom│",
        "│ ain  │ ain  │   │",
        "├──────┴──────┴───┤",
        "│ [Activity Feed] │",
        "└─────────────────┘",
    ]
    add_multiline_box(slide, cx + Inches(0.1), Inches(1.9), Inches(2.6), Inches(2.8),
                      "\n".join(lines), 9, GRAY)

# Bottom description
add_text_box(slide, M, Inches(5.2), Inches(11), Inches(1.5),
    "The platform features a responsive, modern UI with: Dark theme optimized for data dashboards · "
    "Left sidebar navigation with 5 main sections · Search bar and notification system · "
    "Domain-specific filtering across all views · Interactive charts with drill-down · "
    "Mobile-responsive design for on-the-go access", 12, GRAY)

# ============================================================
# SLIDE 10: TECHNOLOGIES USED
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "TECHNOLOGIES USED", 32, WHITE, True)

techs = [
    ("Google Gemini LLM", "Large Language Model for natural language understanding, context-aware responses, and domain-specific expertise", ACCENT),
    ("Google Vertex AI", "ML model training, deployment, and MLOps pipeline for custom predictive models", ACCENT),
    ("Google BigQuery", "Serverless data warehouse for large-scale analytics and time-series queries", ACCENT),
    ("Google Cloud Run", "Fully managed container platform with auto-scaling and pay-per-use pricing", ACCENT),
    ("Next.js 15", "React framework with server-side rendering, app router, and optimal performance", GREEN),
    ("FastAPI (Python)", "High-performance async API framework with automatic OpenAPI documentation", GREEN),
    ("Tailwind CSS", "Utility-first CSS framework for rapid, responsive UI development", GREEN),
    ("Recharts", "Composable charting library built on React components for data visualization", GREEN),
    ("ChromaDB", "Open-source vector database for RAG pipeline and semantic similarity search", ACCENT2),
    ("Docker", "Container runtime for consistent development and production environments", ACCENT2),
    ("Cloud Build", "CI/CD pipeline for automated testing, building, and deployment", ACCENT2),
    ("Redis", "In-memory data store for caching and task queue management", ACCENT2),
]

for i, (name, desc, color) in enumerate(techs):
    col = i % 2
    row = i // 2
    cx = M + col * Inches(6.2)
    cy = Inches(1.2) + row * Inches(0.95)
    add_shape(slide, cx, cy, Inches(5.9), Inches(0.8), CARD_BG)
    add_text_box(slide, cx + Inches(0.12), cy + Inches(0.05), Inches(2.5), Inches(0.3), f"✦ {name}", 12, color, True)
    add_text_box(slide, cx + Inches(0.12), cy + Inches(0.32), Inches(5.6), Inches(0.4), desc, 10, GRAY)

# ============================================================
# SLIDE 11: SNAPSHOT OF PROTOTYPE
# ============================================================
slide = add_slide()
add_shape(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, M, Inches(0.3), Inches(10), Inches(0.7), "SNAPSHOT OF PROTOTYPE", 32, WHITE, True)

snapshot_text = (
    "The Decision Intelligence Platform is fully functional and deployed. Key prototype highlights:\n\n"
    "1. Dashboard — Real-time display of 8 domain metrics with AI-generated insights and activity feed.\n"
    "   Live KPIs: Data sources, Insights generated, Decisions supported, Community score.\n\n"
    "2. AI Chat — Working conversational interface with domain-specific AI agents. Ask questions like:\n"
    "   'Analyze downtown traffic congestion' or 'Predict energy demand for next month'\n\n"
    "3. Analytics — Interactive domain analytics with trend charts, performance metrics, and forecasts.\n"
    "   Each domain has 7+ tracked metrics with historical data (90+ days).\n\n"
    "4. Decision Engine — Input your decision context and options; the AI evaluates and recommends.\n"
    "   Uses multi-criteria analysis with confidence scoring and pros/cons comparison.\n\n"
    "5. Data Management — Connect new data sources through CSV, JSON, PDF, API, or webhook.\n"
    "   Auto-categorization by domain with status monitoring.\n\n"
    "6. Backend API — 7 RESTful endpoints serving analytics, chat, decisions, and data management.\n"
    "   Full auto-generated docs at /docs endpoint.\n\n"
    "All functionality works offline with sample data. Google Cloud integration enhances capabilities\n"
    "when credentials are configured — enabling real Gemini LLM responses and BigQuery analytics."
)
add_multiline_box(slide, M, Inches(1.0), Inches(11.5), Inches(5.5), snapshot_text, 13, GRAY)

# Bottom bar
add_shape(slide, 0, Inches(7.0), W, Inches(0.5), ACCENT)
add_text_box(slide, M, Inches(7.05), Inches(11), Inches(0.35),
             "GitHub: github.com/kushaljagad73/decision-intelligence-platform  |  Built with Google Cloud ☁️",
             11, WHITE, alignment=PP_ALIGN.CENTER)

prs.save(output_path)
print(f"PPT saved to: {output_path}")
